#!/usr/bin/env python3
"""
Multi-Format Document Extractor for Lunar Regolith Simulants

Extracts simulant data from multiple document formats:
- PDF (spec sheets, research papers, books)
- HTML (saved web pages from ResearchGate, etc.)
- PPTX (PowerPoint presentations)
- DOCX (Word documents)

Handles both single-simulant documents (TDS) and multi-simulant documents (research papers).
"""

import re
import json
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Set
from dataclasses import dataclass, field, asdict
from abc import ABC, abstractmethod
import html
from html.parser import HTMLParser

# PDF support
try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# OCR support
try:
    import pytesseract
    from pdf2image import convert_from_path
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# PPTX support
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# DOCX support
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# BeautifulSoup for HTML (optional, fallback to basic parser)
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False


@dataclass
class SimulantData:
    """Container for extracted simulant data"""
    simulant_id: str = ""
    name: str = ""
    source_file: str = ""

    # Basic info
    type: str = ""  # Mare, Highland, etc.
    series: str = ""
    nasa_fom_score: Optional[float] = None
    reference_material: str = ""

    # Physical properties
    density_min: Optional[float] = None
    density_max: Optional[float] = None
    density_mean: Optional[float] = None
    bulk_density: Optional[float] = None
    particle_size_range: str = ""
    particle_size_median: Optional[float] = None
    particle_size_mean: Optional[float] = None
    aspect_ratio: Optional[float] = None
    circularity: Optional[float] = None
    cohesion_kpa: Optional[float] = None
    friction_angle_deg: Optional[float] = None
    angle_of_repose: Optional[float] = None
    magnetic_susceptibility: str = ""
    specific_gravity: Optional[float] = None
    ph_value: Optional[float] = None

    # Composition
    chemical_composition: Dict[str, float] = field(default_factory=dict)
    mineral_composition: Dict[str, float] = field(default_factory=dict)
    mineral_groups: Dict[str, float] = field(default_factory=dict)

    # Raw materials
    raw_materials: List[str] = field(default_factory=list)

    # Metadata
    extraction_confidence: float = 0.0
    extraction_methods: List[str] = field(default_factory=list)
    notes: List[str] = field(default_factory=list)


class HTMLTextExtractor(HTMLParser):
    """Simple HTML text extractor"""
    def __init__(self):
        super().__init__()
        self.text_parts = []
        self.in_script = False
        self.in_style = False

    def handle_starttag(self, tag, attrs):
        if tag == 'script':
            self.in_script = True
        elif tag == 'style':
            self.in_style = True

    def handle_endtag(self, tag):
        if tag == 'script':
            self.in_script = False
        elif tag == 'style':
            self.in_style = False

    def handle_data(self, data):
        if not self.in_script and not self.in_style:
            text = data.strip()
            if text:
                self.text_parts.append(text)

    def get_text(self):
        return ' '.join(self.text_parts)


class MultiFormatExtractor:
    """Extract simulant data from multiple document formats"""

    # Known simulant names and patterns
    KNOWN_SIMULANTS = [
        'JSC-1', 'JSC-1A', 'JSC-1AF', 'JSC-2A',
        'LHS-1', 'LHS-1D', 'LHS-1E', 'LHS-2', 'LHS-2E',
        'LMS-1', 'LMS-1D', 'LMS-1E', 'LMS-2',
        'LSP-1', 'LSP-2',
        'TLH-0', 'TLM-0',
        'EAC-1', 'EAC-1A',
        'FJS-1', 'FJS-2', 'FJS-3',
        'MLS-1', 'MLS-2',
        'NU-LHT-1M', 'NU-LHT-2M', 'NU-LHT-3M',
        'DNA-1', 'DNA-1A',
        'GRC-1', 'GRC-3',
        'BP-1',
        'Chenobi',
        'OPRL2N', 'OPRH2N', 'OPRH3N',
        'CAS-1',
        'CLRS-1', 'CLRS-2',
        'CLDS-1',
        'NAO-1', 'NAO-2',
        'BHLD20',
        'CUG-1A', 'CUG-1B',
        'CUMT-1',
        'AGK-2010',
        'ALRS-1',
        'CSM-CL',
        'KOHLS-1',
        'OB-1',
        'UoM-B', 'UoM-W',
        'NEU-1',
        'TJ-1', 'TJ-2',
        'IGG-01',
        'KLS-1',
        'ISAC-1', 'LSS-ISAC-1',
        'KIGAM',
        'JLRS-1',
    ]

    # Standard oxide patterns
    OXIDE_PATTERNS = {
        'SiO2': [r'SiO[\s\n]*2', r'SiO2', r'Silicon\s*Dioxide'],
        'TiO2': [r'TiO[\s\n]*2', r'TiO2', r'Titanium\s*Dioxide'],
        'Al2O3': [r'Al[\s\n]*2[\s\n]*O[\s\n]*3', r'Al2O3', r'Aluminum\s*Oxide', r'Alumina'],
        'Fe2O3': [r'Fe[\s\n]*2[\s\n]*O[\s\n]*3', r'Fe2O3', r'Ferric\s*Oxide'],
        'FeO': [r'FeO(?![0-9\n])', r'Ferrous\s*Oxide', r'FeO\s+\d'],
        'MgO': [r'MgO', r'Magnesium\s*Oxide'],
        'CaO': [r'CaO', r'Calcium\s*Oxide'],
        'Na2O': [r'Na[\s\n]*2[\s\n]*O', r'Na2O', r'Sodium\s*Oxide'],
        'K2O': [r'K[\s\n]*2[\s\n]*O', r'K2O', r'Potassium\s*Oxide'],
        'P2O5': [r'P[\s\n]*2[\s\n]*O[\s\n]*5', r'P2O5'],
        'MnO': [r'MnO', r'Manganese\s*Oxide'],
        'Cr2O3': [r'Cr[\s\n]*2[\s\n]*O[\s\n]*3', r'Cr2O3'],
        'NiO': [r'NiO', r'Nickel\s*Oxide'],
        'ZnO': [r'ZnO', r'Zinc\s*Oxide'],
        'SrO': [r'SrO', r'Strontium\s*Oxide'],
        'BaO': [r'BaO', r'Barium\s*Oxide'],
        'SO3': [r'SO[\s\n]*3', r'SO3'],
        'LOI': [r'LOI', r'Loss\s*on\s*Ignition'],
    }

    # Mineral patterns (expanded)
    MINERAL_PATTERNS = {
        'plagioclase': [r'plagioclase', r'plag\.?'],
        'anorthite': [r'anorthite'],
        'anorthosite': [r'anorthosite'],
        'pyroxene': [r'pyroxene', r'pyx\.?'],
        'augite': [r'augite'],
        'bronzite': [r'bronzite'],
        'enstatite': [r'enstatite'],
        'clinopyroxene': [r'clinopyroxene', r'cpx'],
        'orthopyroxene': [r'orthopyroxene', r'opx'],
        'olivine': [r'olivine', r'oliv\.?'],
        'forsterite': [r'forsterite', r'fosterite'],
        'fayalite': [r'fayalite'],
        'ilmenite': [r'ilmenite', r'ilm\.?'],
        'glass': [r'\bglass\b', r'glass-rich'],
        'basalt': [r'basalt'],
        'quartz': [r'quartz'],
        'feldspar': [r'feldspar'],
        'magnetite': [r'magnetite'],
        'hematite': [r'hematite'],
        'hornblende': [r'hornblende'],
        'analcime': [r'analcime'],
        'smectite': [r'smectite'],
        'illite': [r'illite'],
        'lizardite': [r'lizardite'],
        'serpentine': [r'serpentine'],
        'agglutinate': [r'agglutinate'],
        'spinel': [r'spinel'],
        'chromite': [r'chromite'],
        'apatite': [r'apatite'],
        'norite': [r'norite'],
        'troctolite': [r'troctolite'],
    }

    def __init__(self):
        self.debug = False
        self.ocr_available = OCR_AVAILABLE
        self.pdf_available = PDF_AVAILABLE
        self.pptx_available = PPTX_AVAILABLE
        self.docx_available = DOCX_AVAILABLE

    def extract_from_file(self, file_path: Path) -> List[SimulantData]:
        """
        Extract simulant data from any supported file format.
        Returns a list of SimulantData (may contain multiple simulants from research papers).
        """
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()

        if suffix == '.pdf':
            return self._extract_from_pdf(file_path)
        elif suffix == '.html' or suffix == '.htm':
            return self._extract_from_html(file_path)
        elif suffix == '.pptx':
            return self._extract_from_pptx(file_path)
        elif suffix in ['.docx', '.doc']:
            return self._extract_from_docx(file_path)
        else:
            return []

    def _extract_from_pdf(self, pdf_path: Path) -> List[SimulantData]:
        """Extract from PDF (spec sheets and research papers)"""
        if not self.pdf_available:
            return []

        results = []
        full_text = ""
        all_tables = []

        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    full_text += text + "\n"
                    tables = page.extract_tables()
                    if tables:
                        all_tables.extend(tables)

            # Detect simulants mentioned in the document
            mentioned_simulants = self._find_mentioned_simulants(full_text)

            # Check if this is a spec sheet (single simulant focus)
            is_spec_sheet = self._is_spec_sheet(pdf_path.name, full_text)

            if is_spec_sheet:
                # Single simulant extraction
                data = self._extract_single_simulant(pdf_path.name, full_text, all_tables)
                if data.name:
                    results.append(data)
            else:
                # Research paper - try to extract data for each mentioned simulant
                for simulant_name in mentioned_simulants:
                    data = self._extract_simulant_from_paper(
                        simulant_name, pdf_path.name, full_text, all_tables
                    )
                    if data and (data.chemical_composition or data.mineral_composition or data.type):
                        results.append(data)

        except Exception as e:
            if self.debug:
                print(f"PDF extraction error: {e}")

        return results

    def _extract_from_html(self, html_path: Path) -> List[SimulantData]:
        """Extract from HTML files (saved web pages)"""
        results = []

        try:
            with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            # Extract text
            if BS4_AVAILABLE:
                soup = BeautifulSoup(content, 'html.parser')
                # Remove scripts and styles
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text(separator=' ', strip=True)
            else:
                parser = HTMLTextExtractor()
                parser.feed(content)
                text = parser.get_text()

            # Find mentioned simulants
            mentioned_simulants = self._find_mentioned_simulants(text)

            for simulant_name in mentioned_simulants:
                data = self._extract_simulant_from_paper(
                    simulant_name, html_path.name, text, []
                )
                if data and (data.chemical_composition or data.mineral_composition or data.type):
                    data.extraction_methods.append('html')
                    results.append(data)

        except Exception as e:
            if self.debug:
                print(f"HTML extraction error: {e}")

        return results

    def _extract_from_pptx(self, pptx_path: Path) -> List[SimulantData]:
        """Extract from PowerPoint files"""
        if not self.pptx_available:
            return []

        results = []
        full_text = ""
        tables_data = []

        try:
            prs = Presentation(pptx_path)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"

                    # Extract table data
                    if shape.has_table:
                        table = shape.table
                        table_rows = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_rows.append(row_data)
                        tables_data.append(table_rows)

            # Detect simulants
            mentioned_simulants = self._find_mentioned_simulants(full_text)

            if self._is_spec_sheet(pptx_path.name, full_text):
                data = self._extract_single_simulant(pptx_path.name, full_text, tables_data)
                if data.name:
                    data.extraction_methods.append('pptx')
                    results.append(data)
            else:
                for simulant_name in mentioned_simulants:
                    data = self._extract_simulant_from_paper(
                        simulant_name, pptx_path.name, full_text, tables_data
                    )
                    if data:
                        data.extraction_methods.append('pptx')
                        results.append(data)

        except Exception as e:
            if self.debug:
                print(f"PPTX extraction error: {e}")

        return results

    def _extract_from_docx(self, docx_path: Path) -> List[SimulantData]:
        """Extract from Word documents"""
        if not self.docx_available:
            return []

        results = []
        full_text = ""
        tables_data = []

        try:
            doc = Document(docx_path)

            # Extract paragraphs
            for para in doc.paragraphs:
                full_text += para.text + "\n"

            # Extract tables
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_rows.append(row_data)
                tables_data.append(table_rows)

            mentioned_simulants = self._find_mentioned_simulants(full_text)

            for simulant_name in mentioned_simulants:
                data = self._extract_simulant_from_paper(
                    simulant_name, docx_path.name, full_text, tables_data
                )
                if data:
                    data.extraction_methods.append('docx')
                    results.append(data)

        except Exception as e:
            if self.debug:
                print(f"DOCX extraction error: {e}")

        return results

    def _find_mentioned_simulants(self, text: str) -> List[str]:
        """Find all simulant names mentioned in text"""
        found = []
        text_upper = text.upper()

        for simulant in self.KNOWN_SIMULANTS:
            # Create pattern that matches the simulant name with word boundaries
            pattern = r'\b' + re.escape(simulant.upper()) + r'\b'
            if re.search(pattern, text_upper):
                found.append(simulant)

        return found

    def _is_spec_sheet(self, filename: str, text: str) -> bool:
        """Determine if document is a spec sheet (single simulant focus)"""
        filename_lower = filename.lower()

        # Check filename patterns
        spec_patterns = ['tds', 'spec', 'sheet', 'datasheet', 'fact']
        if any(p in filename_lower for p in spec_patterns):
            return True

        # Check content patterns
        text_lower = text.lower()[:2000]
        if any(p in text_lower for p in ['technical data sheet', 'fact sheet', 'specification sheet']):
            return True

        # Check if only one simulant is prominently featured
        mentioned = self._find_mentioned_simulants(text)
        if len(mentioned) == 1:
            return True

        return False

    def _extract_single_simulant(self, filename: str, text: str, tables: List) -> SimulantData:
        """Extract data for a single simulant from spec sheet"""
        data = SimulantData()
        data.source_file = filename

        # Extract name from filename or text
        data.name = self._extract_name(filename, text)

        # Extract all properties
        self._extract_basic_info(text, data)
        self._extract_chemical_composition(text, tables, data)
        self._extract_mineral_composition(text, tables, data)
        self._extract_physical_properties(text, data)

        data.extraction_confidence = self._calculate_confidence(data)
        return data

    def _extract_simulant_from_paper(
        self, simulant_name: str, filename: str, text: str, tables: List
    ) -> Optional[SimulantData]:
        """Extract data for a specific simulant from a research paper"""
        data = SimulantData()
        data.name = simulant_name
        data.source_file = filename

        # Find context around simulant mentions
        context = self._get_simulant_context(simulant_name, text)

        # Extract properties from context
        self._extract_basic_info(context, data)
        self._extract_chemical_from_context(simulant_name, text, tables, data)
        self._extract_mineral_from_context(simulant_name, text, tables, data)
        self._extract_physical_properties(context, data)

        # Only return if we found meaningful data
        if data.chemical_composition or data.mineral_composition or data.type or data.nasa_fom_score:
            data.extraction_methods.append('research_paper')
            data.extraction_confidence = self._calculate_confidence(data)
            return data

        return None

    def _get_simulant_context(self, simulant_name: str, text: str, window: int = 1500) -> str:
        """Get text context around simulant mentions"""
        contexts = []
        pattern = re.compile(re.escape(simulant_name), re.IGNORECASE)

        for match in pattern.finditer(text):
            start = max(0, match.start() - window)
            end = min(len(text), match.end() + window)
            contexts.append(text[start:end])

        return '\n'.join(contexts)

    def _extract_name(self, filename: str, text: str) -> str:
        """Extract simulant name from filename or text"""
        # Try filename first
        name = filename.replace('.pdf', '').replace('.html', '').replace('.pptx', '').replace('.docx', '')

        # Check for known simulant patterns in filename
        for simulant in self.KNOWN_SIMULANTS:
            if simulant.upper() in name.upper():
                return simulant

        # Common patterns
        patterns = [
            r'^([A-Z]{2,3}-\d+[A-Z]?)',  # LMS-1, LHS-1D
            r'([A-Z]{2,3}-\d+[A-Z]?)[_\s-]',
            r'Simulant[:\s]+([A-Z0-9-]+)',
        ]

        for pattern in patterns:
            match = re.search(pattern, name, re.IGNORECASE)
            if match:
                return match.group(1).upper()

        # Check text
        text_start = text[:500]
        for simulant in self.KNOWN_SIMULANTS:
            if simulant.upper() in text_start.upper():
                return simulant

        return ""

    def _extract_basic_info(self, text: str, data: SimulantData):
        """Extract basic simulant info"""
        # Type
        type_patterns = [
            r'(?:Simulant\s+)?Type[:\s]+([A-Za-z]+(?:[-\s][A-Za-z]+)*)',
            r'((?:Low-Ti\s+)?(?:High-Ti\s+)?(?:Lunar\s+)?(?:Mare|Highland|Maria))',
            r'(Mare|Highland)\s+(?:Simulant|Type)',
        ]
        for pattern in type_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                type_val = match.group(1).strip()
                if len(type_val) < 30:  # Sanity check
                    data.type = type_val
                    break

        # Reference material
        ref_match = re.search(r'Reference\s+Material[:\s]+([^\n]+)', text, re.IGNORECASE)
        if ref_match:
            data.reference_material = ref_match.group(1).strip()[:100]

        # NASA FoM Score
        fom_patterns = [
            r'(?:NASA\s+)?FoM\s*(?:Score)?[:\s]*(\d+\.?\d*)\s*%?',
            r'Figure[s]?\s+of\s+Merit[:\s]*(\d+\.?\d*)\s*%?',
        ]
        for pattern in fom_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                try:
                    fom = float(match.group(1))
                    if 0 <= fom <= 100:
                        data.nasa_fom_score = fom
                        break
                except ValueError:
                    pass

    def _extract_chemical_composition(self, text: str, tables: List, data: SimulantData):
        """Extract chemical composition from text and tables"""
        # Method 1: Parse inline text for oxide values
        self._parse_oxide_text(text, data)

        # Method 2: Parse tables
        if len(data.chemical_composition) < 5:
            self._parse_oxide_tables(tables, data)

        # Method 3: Parse formatted oxide blocks
        if len(data.chemical_composition) < 5:
            self._parse_oxide_block(text, data)

    def _extract_chemical_from_context(self, simulant_name: str, text: str, tables: List, data: SimulantData):
        """Extract chemical composition for a specific simulant"""
        # Find tables that mention this simulant
        for table in tables:
            if not table:
                continue

            table_text = str(table)
            if simulant_name.upper() in table_text.upper():
                self._parse_oxide_from_table(table, simulant_name, data)

        # Parse from context text
        context = self._get_simulant_context(simulant_name, text, window=2000)
        self._parse_oxide_text(context, data)

    def _parse_oxide_text(self, text: str, data: SimulantData):
        """Parse oxide values from text"""
        for oxide, patterns in self.OXIDE_PATTERNS.items():
            if oxide in data.chemical_composition:
                continue

            for pattern in patterns:
                # Look for oxide followed by value
                regex = rf'{pattern}[:\s]*(\d+\.?\d*)\s*(?:%|wt\.?%?)?'
                matches = re.findall(regex, text, re.IGNORECASE)
                for match in matches:
                    try:
                        value = float(match)
                        if 0 <= value <= 100:
                            data.chemical_composition[oxide] = value
                            break
                    except (ValueError, IndexError):
                        pass

    def _parse_oxide_block(self, text: str, data: SimulantData):
        """Parse oxide composition from formatted blocks"""
        # Pattern: Oxide Wt.% or similar headers followed by values
        lines = text.split('\n')

        for i, line in enumerate(lines):
            # Look for oxide header patterns
            if re.search(r'Oxide\s+(?:Wt\.?%|Weight)', line, re.IGNORECASE):
                # Parse following lines
                for j in range(i + 1, min(i + 20, len(lines))):
                    oxide_line = lines[j]
                    # Match patterns like "SiO2 48.22" or "SiO 48.22"
                    for oxide, patterns in self.OXIDE_PATTERNS.items():
                        if oxide in data.chemical_composition:
                            continue
                        for pattern in patterns:
                            match = re.search(rf'({pattern})\s+(\d+\.?\d*)', oxide_line, re.IGNORECASE)
                            if match:
                                try:
                                    value = float(match.group(2))
                                    if 0 <= value <= 100:
                                        data.chemical_composition[oxide] = value
                                except ValueError:
                                    pass

    def _parse_oxide_tables(self, tables: List, data: SimulantData):
        """Parse oxide composition from extracted tables"""
        for table in tables:
            if not table or len(table) < 2:
                continue

            # Look for oxide headers in first row
            header_row = [str(cell).strip() if cell else '' for cell in table[0]]
            header_text = ' '.join(header_row).upper()

            if any(ox in header_text for ox in ['SIO', 'TIO', 'AL2O', 'FEO', 'MGO', 'CAO']):
                self._parse_oxide_from_table(table, None, data)

    def _parse_oxide_from_table(self, table: List, simulant_name: Optional[str], data: SimulantData):
        """Parse oxide values from a specific table"""
        if not table or len(table) < 2:
            return

        # Try to identify header row and value row
        for row_idx, row in enumerate(table):
            row_text = ' '.join(str(cell) if cell else '' for cell in row)

            # Check if this row contains oxide names
            oxide_count = sum(1 for ox in ['SiO', 'TiO', 'Al2O', 'Al O', 'FeO', 'MgO', 'CaO']
                            if ox in row_text)

            if oxide_count >= 3:
                # This is likely a header row - find corresponding value row
                for val_idx in range(row_idx + 1, min(row_idx + 3, len(table))):
                    value_row = table[val_idx]
                    if simulant_name:
                        # Check if this row is for our simulant
                        row_str = str(value_row)
                        if simulant_name.upper() not in row_str.upper():
                            continue

                    # Extract values
                    for col_idx, cell in enumerate(row):
                        if not cell:
                            continue
                        cell_str = str(cell).strip()

                        # Match cell to oxide
                        for oxide, patterns in self.OXIDE_PATTERNS.items():
                            if oxide in data.chemical_composition:
                                continue
                            for pattern in patterns:
                                if re.search(pattern, cell_str, re.IGNORECASE):
                                    # Get value from corresponding column in value row
                                    if col_idx < len(value_row) and value_row[col_idx]:
                                        try:
                                            val_str = str(value_row[col_idx]).strip()
                                            value = float(re.sub(r'[^\d.]', '', val_str))
                                            if 0 <= value <= 100:
                                                data.chemical_composition[oxide] = value
                                        except (ValueError, IndexError):
                                            pass
                                    break

    def _extract_mineral_composition(self, text: str, tables: List, data: SimulantData):
        """Extract mineral composition"""
        self._parse_mineral_text(text, data)

        if len(data.mineral_composition) < 3:
            self._parse_mineral_tables(tables, data)

    def _extract_mineral_from_context(self, simulant_name: str, text: str, tables: List, data: SimulantData):
        """Extract mineral composition for specific simulant"""
        context = self._get_simulant_context(simulant_name, text, window=2000)
        self._parse_mineral_text(context, data)

        for table in tables:
            if not table:
                continue
            table_text = str(table)
            if simulant_name.upper() in table_text.upper():
                self._parse_mineral_from_table(table, simulant_name, data)

    def _parse_mineral_text(self, text: str, data: SimulantData):
        """Parse mineral values from text"""
        for mineral, patterns in self.MINERAL_PATTERNS.items():
            if mineral in data.mineral_composition:
                continue

            for pattern in patterns:
                regex = rf'{pattern}[:\s]*(\d+\.?\d*)\s*(?:%|wt\.?%?)?'
                matches = re.findall(regex, text, re.IGNORECASE)
                for match in matches:
                    try:
                        value = float(match)
                        if 0 <= value <= 100:
                            data.mineral_composition[mineral] = value
                            break
                    except (ValueError, IndexError):
                        pass

    def _parse_mineral_tables(self, tables: List, data: SimulantData):
        """Parse mineral composition from tables"""
        for table in tables:
            if not table or len(table) < 2:
                continue

            header_row = [str(cell).strip() if cell else '' for cell in table[0]]
            header_text = ' '.join(header_row).lower()

            mineral_keywords = ['plagioclase', 'pyroxene', 'olivine', 'glass', 'ilmenite', 'mineral']
            if any(kw in header_text for kw in mineral_keywords):
                self._parse_mineral_from_table(table, None, data)

    def _parse_mineral_from_table(self, table: List, simulant_name: Optional[str], data: SimulantData):
        """Parse mineral values from table"""
        if not table or len(table) < 2:
            return

        for row in table:
            row_text = ' '.join(str(cell) if cell else '' for cell in row).lower()

            for mineral, patterns in self.MINERAL_PATTERNS.items():
                if mineral in data.mineral_composition:
                    continue

                for pattern in patterns:
                    if re.search(pattern, row_text, re.IGNORECASE):
                        # Find numeric value in this row
                        for cell in row:
                            if cell:
                                try:
                                    val_str = str(cell).strip()
                                    # Skip if this is the mineral name itself
                                    if re.search(pattern, val_str, re.IGNORECASE):
                                        continue
                                    value = float(re.sub(r'[^\d.]', '', val_str))
                                    if 0 < value <= 100:
                                        data.mineral_composition[mineral] = value
                                        break
                                except (ValueError, IndexError):
                                    pass

    def _extract_physical_properties(self, text: str, data: SimulantData):
        """Extract physical properties"""
        # Bulk Density
        density_patterns = [
            r'(?:Bulk|Uncompressed)\s+Density[:\s]*(\d+\.?\d*)\s*g/cm',
            r'Density[:\s]*(\d+\.?\d*)\s*g/cm',
            r'(\d+\.?\d*)\s*g/cm[³3]?\s*(?:bulk|density)',
        ]
        for pattern in density_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                try:
                    value = float(match.group(1))
                    if 0.5 <= value <= 5:
                        data.bulk_density = value
                        break
                except ValueError:
                    pass

        # Median Particle Size
        size_patterns = [
            r'Median\s+Particle\s+Size[:\s]*(\d+\.?\d*)\s*[µμu]m',
            r'D50[:\s=]*(\d+\.?\d*)\s*[µμu]m',
            r'(\d+\.?\d*)\s*[µμu]m\s*(?:median|D50)',
        ]
        for pattern in size_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                try:
                    value = float(match.group(1))
                    if 0.1 <= value <= 10000:
                        data.particle_size_median = value
                        break
                except ValueError:
                    pass

        # Cohesion
        cohesion_match = re.search(r'Cohesion[:\s]*(\d+\.?\d*)\s*kPa', text, re.IGNORECASE)
        if cohesion_match:
            try:
                data.cohesion_kpa = float(cohesion_match.group(1))
            except ValueError:
                pass

        # Friction Angle
        friction_patterns = [
            r'(?:Angle\s+of\s+)?(?:Internal\s+)?Friction[:\s]*(\d+\.?\d*)[°\s]',
            r'Friction\s+Angle[:\s]*(\d+\.?\d*)',
        ]
        for pattern in friction_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                try:
                    value = float(match.group(1))
                    if 0 <= value <= 90:
                        data.friction_angle_deg = value
                        break
                except ValueError:
                    pass

        # pH
        ph_match = re.search(r'pH[:\s]*(\d+\.?\d*)', text, re.IGNORECASE)
        if ph_match:
            try:
                value = float(ph_match.group(1))
                if 0 <= value <= 14:
                    data.ph_value = value
            except ValueError:
                pass

        # Specific Gravity
        sg_match = re.search(r'Specific\s+Gravity[:\s]*(\d+\.?\d*)', text, re.IGNORECASE)
        if sg_match:
            try:
                value = float(sg_match.group(1))
                if 1 <= value <= 5:
                    data.specific_gravity = value
            except ValueError:
                pass

    def _calculate_confidence(self, data: SimulantData) -> float:
        """Calculate extraction confidence score"""
        score = 0.0

        if data.name:
            score += 20

        if data.type:
            score += 10

        if len(data.chemical_composition) >= 5:
            score += 25
        elif len(data.chemical_composition) >= 3:
            score += 15
        elif len(data.chemical_composition) >= 1:
            score += 5

        if len(data.mineral_composition) >= 3:
            score += 20
        elif len(data.mineral_composition) >= 1:
            score += 10

        if data.bulk_density:
            score += 5
        if data.particle_size_median:
            score += 5
        if data.nasa_fom_score:
            score += 10
        if data.cohesion_kpa:
            score += 5

        return min(score, 100)


def main():
    """Test the multi-format extractor"""
    extractor = MultiFormatExtractor()
    extractor.debug = True

    print("Multi-Format Extractor Test")
    print("=" * 60)
    print(f"PDF support: {extractor.pdf_available}")
    print(f"OCR support: {extractor.ocr_available}")
    print(f"PPTX support: {extractor.pptx_available}")
    print(f"DOCX support: {extractor.docx_available}")
    print()

    # Test with files from DIRT Papers
    test_dir = Path("/home/alvaro/Spring - Forest on the moon/DIRT/DIRT Papers/")

    if test_dir.exists():
        files = list(test_dir.glob("*.pdf"))[:5]
        files += list(test_dir.glob("*.html"))[:3]

        for file_path in files:
            print(f"\n{'=' * 60}")
            print(f"Processing: {file_path.name}")
            print("=" * 60)

            results = extractor.extract_from_file(file_path)

            if results:
                for data in results:
                    print(f"\n  Simulant: {data.name}")
                    print(f"  Type: {data.type or 'N/A'}")
                    print(f"  Confidence: {data.extraction_confidence}%")
                    print(f"  Chemicals: {len(data.chemical_composition)}")
                    print(f"  Minerals: {len(data.mineral_composition)}")
                    if data.chemical_composition:
                        print(f"  Sample oxides: {list(data.chemical_composition.items())[:3]}")
            else:
                print("  No simulant data extracted")


if __name__ == "__main__":
    main()
